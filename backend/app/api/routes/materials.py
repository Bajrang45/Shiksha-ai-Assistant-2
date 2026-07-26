import re
from io import BytesIO

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile, status

from app.api.dependencies import get_current_user
from app.schemas.materials import MaterialResponse
from app.services.material_store import material_store

router = APIRouter(prefix="/materials", tags=["Materials"])
SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".webp"}


def short_summary(text: str) -> str:
    sentences = re.split(r"(?<=[.!?])\s+", re.sub(r"\s+", " ", text).strip())
    useful = [sentence.strip() for sentence in sentences if len(sentence.strip()) > 35]
    selection = useful[:3] or sentences[:3]
    return " ".join(selection)[:900] or "The file did not contain readable text."


def extract_text(filename: str, content: bytes) -> str:
    extension = filename.lower().rsplit(".", 1)[-1]
    if extension == "txt":
        return content.decode("utf-8", errors="replace")
    if extension == "pdf":
        try:
            import fitz
            document = fitz.open(stream=content, filetype="pdf")
            return "\n".join(f"[[PAGE {index + 1}]]\n{page.get_text()}" for index, page in enumerate(document))
        except ImportError as exc:
            raise HTTPException(status_code=503, detail="PDF support is not installed. Run pip install -r requirements.txt.") from exc
    if extension == "docx":
        try:
            from docx import Document
            return "\n".join(paragraph.text for paragraph in Document(BytesIO(content)).paragraphs)
        except ImportError as exc:
            raise HTTPException(status_code=503, detail="DOCX support is not installed. Run pip install -r requirements.txt.") from exc
    if extension == "pptx":
        try:
            from pptx import Presentation
            presentation = Presentation(BytesIO(content))
            return "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        except ImportError as exc:
            raise HTTPException(status_code=503, detail="PPTX support is not installed. Run pip install -r requirements.txt.") from exc
    if extension in {"png", "jpg", "jpeg", "webp"}:
        try:
            import easyocr
            from PIL import Image
            import numpy as np
            image = np.array(Image.open(BytesIO(content)).convert("RGB"))
            return "\n".join(easyocr.Reader(["en"], gpu=False).readtext(image, detail=0))
        except ImportError as exc:
            raise HTTPException(status_code=503, detail="Image text recognition is not installed. Run pip install -r requirements.txt.") from exc
    raise HTTPException(status_code=415, detail="Supported files: PDF, DOCX, PPTX, TXT, and images.")


@router.post("/upload", response_model=MaterialResponse, status_code=status.HTTP_201_CREATED)
async def upload_material(file: UploadFile = File(...), current_user: dict = Depends(get_current_user)) -> MaterialResponse:
    filename = file.filename or "study-material"
    if not any(filename.lower().endswith(extension) for extension in SUPPORTED_EXTENSIONS):
        raise HTTPException(status_code=415, detail="Supported files: PDF, DOCX, PPTX, TXT, and images.")
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="The uploaded file is empty.")
    if len(content) > 10 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Files must be 10 MB or smaller.")
    text = extract_text(filename, content).strip()
    if not text:
        raise HTTPException(status_code=422, detail="No readable text was found in this file. Try a text-based PDF or DOCX.")
    await material_store.add(current_user["id"], filename, text)
    return MaterialResponse(filename=filename, summary=short_summary(text), characters_extracted=len(text))
